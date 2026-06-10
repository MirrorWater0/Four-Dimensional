from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_DIR = Path("docs/presentation")
OUTPUT_FILE = OUTPUT_DIR / "Career_Survey_Presentation_Draft.pptx"
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

PRIMARY = RGBColor(35, 70, 120)
ACCENT = RGBColor(232, 120, 67)
BG = RGBColor(246, 248, 252)
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(92, 102, 112)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(8.7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = PRIMARY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.45), Inches(8.5), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED


def add_bullets(slide, left, top, width, height, lines, font_size=20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_note_box(slide, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(0.82), Inches(2.1), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_bar_chart(slide, left, top, width, height, categories, values, title):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Students", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.value_axis.maximum_scale = max(values) + 2
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.has_legend = False

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PRIMARY

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(11)


def add_pie_chart(slide, left, top, width, height, categories, values, title):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Share", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = True
    if chart.legend is not None:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_legend_key = False
    plot.data_labels.font.size = Pt(11)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.62), Inches(6.78), Inches(8.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = MUTED


def build_presentation():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Career Choices Among College Students", "A short oral presentation draft for a 3-minute survey report")
    add_note_box(slide, "Editable Draft")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(7.4),
        Inches(3.0),
        [
            "Topic: the most important factors when students choose a future career",
            "Reason: career decisions affect motivation, study plans, and long-term life goals",
            "This draft follows your teacher's required structure: topic, survey, findings, discussion, and suggestions",
        ],
        font_size=20,
    )
    add_footer(slide, "Replace the sample numbers with your own survey results before presenting.")

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Survey Design and Participants")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.8),
        Inches(4.5),
        Inches(3.8),
        [
            "Questionnaire survey in English or Chinese",
            "Target group: college students from different majors",
            "Recommended sample: at least 20 students",
            "Suggested background data: major, year of study, and gender ratio",
            "Core questions: 6 to 8 career-related items",
        ],
        font_size=18,
    )

    info = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(1.95), Inches(3.65), Inches(3.0))
    info.fill.solid()
    info.fill.fore_color.rgb = RGBColor(255, 255, 255)
    info.line.color.rgb = ACCENT
    tf = info.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(
        [
            "Example participant profile",
            "24 students answered the survey",
            "Majors: Business, IT, Education, Design",
            "Years: 10 freshmen, 7 sophomores, 5 juniors, 2 seniors",
            "Gender ratio: 14 female, 10 male",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(16 if idx == 0 else 14)
        p.font.bold = idx == 0
        p.font.color.rgb = PRIMARY if idx == 0 else TEXT
        p.space_after = Pt(8)
    add_footer(slide, "Keep this slide brief: who answered, how many people, and what questions you asked.")

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Survey Questions")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.8),
        Inches(8.5),
        Inches(4.8),
        [
            "What career field are you most interested in?",
            "Which factor matters most when choosing a career: salary, passion, stability, or family expectations?",
            "Would you prefer a job related to your major?",
            "Which career field seems most promising in the future?",
            "What is your biggest worry about your future job?",
            "What can schools do to help students make career decisions?",
        ],
        font_size=18,
    )
    add_footer(slide, "These questions fit the teacher's topic examples and are easy to explain in 30 to 40 seconds.")

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Key Finding 1: What Matters Most?")
    add_bar_chart(
        slide,
        Inches(0.55),
        Inches(1.75),
        Inches(5.4),
        Inches(3.8),
        ["Salary", "Passion", "Stability", "Family\nexpectations"],
        [6, 10, 7, 1],
        "Most important factor in career choice",
    )
    add_bullets(
        slide,
        Inches(6.2),
        Inches(2.0),
        Inches(3.0),
        Inches(3.5),
        [
            "Passion ranked first",
            "Job stability was the second strongest factor",
            "Family expectations had very limited influence",
        ],
        font_size=16,
    )
    add_footer(slide, "Sample chart language: The bar chart shows that passion was the top factor for most students.")

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Key Finding 2: Future Career Fields")
    add_pie_chart(
        slide,
        Inches(0.75),
        Inches(1.75),
        Inches(4.8),
        Inches(3.8),
        ["Technology", "Healthcare", "Education", "Business", "Design"],
        [9, 6, 4, 3, 2],
        "Most promising fields for the future",
    )
    add_bullets(
        slide,
        Inches(5.9),
        Inches(1.95),
        Inches(3.1),
        Inches(3.6),
        [
            "Technology was the most popular choice",
            "Healthcare also received strong support",
            "Students may connect these fields with high demand and stable income",
        ],
        font_size=16,
    )
    add_footer(slide, "You can say: Overall, students believed technology and healthcare offer the best opportunities.")

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    add_title(slide, "Discussion, Suggestions, and Conclusion")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.8),
        Inches(8.6),
        Inches(4.6),
        [
            "The findings suggest that students want both personal interest and practical security.",
            "Many students are optimistic about fast-growing industries such as technology and healthcare.",
            "Schools should provide more career talks, internship information, and one-to-one guidance.",
            "Students should explore careers early and compare their interests with real job market needs.",
        ],
        font_size=18,
    )
    add_footer(slide, "End with one sentence for Q&A: Thank you for listening. I am ready for your questions.")

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
